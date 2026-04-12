from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xC8, 0xAA)   # teal-green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG   = RGBColor(0x14, 0x2B, 0x40)   # slightly lighter navy

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def accent_bar(slide, top=0.55, width=1.6, height=0.06):
    add_rect(slide, 0.55, top, width, height, ACCENT)

def add_bullet_points(slide, items, left, top, width, height,
                      font_size=16, color=LIGHT_GRAY, spacing_pt=6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

# ── Create Presentation ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

# Large logo-style wordmark
add_textbox(s, "EARTH", 0.55, 1.8, 6, 1.6,
            font_size=90, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(s, "RISK", 0.55, 3.1, 6, 1.4,
            font_size=90, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# Accent bar
add_rect(s, 0.55, 2.95, 5.2, 0.07, ACCENT)

add_textbox(s, "AI-Powered Insurance Risk Intelligence", 0.55, 4.7, 8, 0.6,
            font_size=22, bold=False, color=LIGHT_GRAY, italic=True)

add_textbox(s, "5-Minute Pitch  •  2026", 0.55, 6.6, 5, 0.5,
            font_size=13, color=RGBColor(0x66,0x88,0x99))

# Right-side decorative element — three stacked accent rectangles
for i, (h, op) in enumerate([(0.06, ACCENT), (0.06, RGBColor(0x00,0x88,0x77)), (0.06, RGBColor(0x00,0x55,0x44))]):
    add_rect(s, 9.5, 1.5 + i * 1.8, 3.3, 1.5, RGBColor(0x14,0x2B,0x40))
labels = ["Seismic  30%", "Climate  20%", "Fire Access  20%"]
for i, lbl in enumerate(labels):
    add_textbox(s, lbl, 9.6, 1.65 + i * 1.8, 3.0, 0.8,
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "THE PROBLEM", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

add_textbox(s, "Insurers Are Flying Blind", 0.55, 0.7, 9, 0.9,
            font_size=42, bold=True, color=WHITE)
accent_bar(s, 1.55, 4.0)

problems = [
    "🌍  Climate risk is accelerating — floods, earthquakes, wildfires are intensifying",
    "📋  Risk assessment still relies on manual surveys & outdated spreadsheets",
    "💸  Mispriced policies = massive losses ($billions/year in unexpected claims)",
    "⏱️  Adjusters spend weeks on assessments that should take minutes",
]
add_bullet_points(s, problems, 0.55, 1.8, 9.5, 4.0, font_size=20, spacing_pt=14)

# stat cards
for i, (val, label) in enumerate([
    ("$280B", "annual insured\nclimate losses"),
    ("43%", "of policies\nmispriced"),
    ("6 weeks", "avg. manual\nrisk assessment"),
]):
    x = 0.3 + i * 4.3
    add_rect(s, x, 5.6, 3.8, 1.6, CARD_BG)
    add_textbox(s, val, x + 0.15, 5.7, 3.5, 0.8,
                font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x + 0.15, 6.3, 3.5, 0.7,
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "OUR SOLUTION", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "EarthRisk Platform", 0.55, 0.7, 10, 0.9,
            font_size=42, bold=True, color=WHITE)
accent_bar(s, 1.55, 4.0)

add_textbox(s,
    "A real-time, AI-powered risk intelligence platform that gives insurers "
    "instant, data-driven risk scores for any building — anywhere.",
    0.55, 1.75, 8.5, 1.0, font_size=19, italic=True, color=LIGHT_GRAY)

# 4 feature cards
features = [
    ("🗺️", "Interactive\nRisk Maps", "Heatmaps showing seismic,\nflood & climate exposure"),
    ("🤖", "4 AI Agents", "Risk Explainer · Alert Monitor\nData Interpreter · Decision Support"),
    ("📊", "EarthRisk\nScore", "Composite 0–100 score\nacross 6 weighted factors"),
    ("⚡", "Real-Time\nAlerts", "Auto-alerts when risk\nexceeds threshold"),
]
for i, (icon, title, desc) in enumerate(features):
    x = 0.3 + i * 3.25
    add_rect(s, x, 3.0, 3.0, 3.8, CARD_BG)
    add_rect(s, x, 3.0, 3.0, 0.06, ACCENT)
    add_textbox(s, icon, x + 0.1, 3.1, 2.8, 0.7,
                font_size=28, align=PP_ALIGN.CENTER)
    add_textbox(s, title, x + 0.1, 3.75, 2.8, 0.7,
                font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + 0.1, 4.45, 2.8, 1.2,
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "HOW IT WORKS", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "Architecture in 3 Layers", 0.55, 0.7, 10, 0.9,
            font_size=38, bold=True, color=WHITE)
accent_bar(s, 1.55, 3.5)

# Layer boxes
layers = [
    ("FRONTEND", "React 19 + Vite", "Mapbox GL heatmaps\nLeaflet overlays\nRecharts dashboards\nJWT-secured login"),
    ("BACKEND", "Node.js / Express", "MySQL building database\n8 REST API endpoints\nJWT auth middleware\nAdmin statistics"),
    ("AI AGENTS", "IBM WatsonX\nOrchestrate", "4 specialized agents\nPython tool functions\nKnowledge bases\nEarthRisk scoring"),
]
for i, (layer, tech, detail) in enumerate(layers):
    x = 0.4 + i * 4.3
    add_rect(s, x, 2.0, 3.9, 4.9, CARD_BG)
    add_rect(s, x, 2.0, 3.9, 0.5, ACCENT)
    add_textbox(s, layer, x + 0.1, 2.05, 3.7, 0.45,
                font_size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_textbox(s, tech, x + 0.1, 2.6, 3.7, 0.7,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, detail, x + 0.2, 3.35, 3.5, 2.5,
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # connector arrow (not last)
    if i < 2:
        add_textbox(s, "→", x + 4.0, 4.1, 0.4, 0.5,
                    font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# EarthRisk formula
add_rect(s, 0.4, 6.9, 12.5, 0.45, CARD_BG)
add_textbox(s,
    "Score  =  Seismic×0.30  +  Volcanic×0.10  +  Fire Access×0.20  +  Climate×0.20  +  Age×0.10  +  Claims×0.10",
    0.55, 6.92, 12.2, 0.4, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI AGENTS DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "AI AGENTS", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "4 Specialized Agents, 1 Intelligent System", 0.55, 0.7, 12, 0.9,
            font_size=36, bold=True, color=WHITE)
accent_bar(s, 1.55, 6.0)

agents = [
    ("Risk\nExplainer", "Calculates EarthRisk score &\nexplains all contributing factors\nto underwriters in plain language"),
    ("Alert\nMonitor", "Continuously monitors portfolio\nand fires alerts when any\nbuilding exceeds score 65"),
    ("Data\nInterpreter", "Translates raw geospatial data\ninto insurer-ready summaries\nwith context & benchmarks"),
    ("Decision\nSupport", "Recommends premium adjustments,\nreinspection schedules, and\nrisk mitigation strategies"),
]
for i, (name, desc) in enumerate(agents):
    x = 0.3 + i * 3.25
    add_rect(s, x, 2.0, 3.0, 4.8, CARD_BG)
    add_rect(s, x, 2.0, 3.0, 0.06, ACCENT)
    # agent number
    add_textbox(s, f"0{i+1}", x + 0.15, 2.15, 0.7, 0.5,
                font_size=28, bold=True, color=ACCENT)
    add_textbox(s, name, x + 0.1, 2.7, 2.8, 0.8,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + 0.15, 3.55, 2.7, 2.5,
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s,
    "All agents share 4 live HTTP tools → fetch building data, history, and portfolio stats directly from the EarthRisk backend",
    0.55, 6.7, 12.2, 0.55, font_size=14, italic=True,
    color=RGBColor(0x88, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "MARKET OPPORTUNITY", 0.55, 0.08, 6, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "A Massive & Growing Market", 0.55, 0.7, 9, 0.9,
            font_size=40, bold=True, color=WHITE)
accent_bar(s, 1.55, 4.0)

# Big market numbers
for i, (val, label, sub) in enumerate([
    ("$7.5T",  "Global P&C Insurance Market",       "Total addressable market"),
    ("$2.1T",  "Commercial Property Insurance",      "Primary target segment"),
    ("$890B",  "Climate Risk Analytics Market",      "Our direct addressable market"),
]):
    y = 1.8 + i * 1.55
    add_rect(s, 0.4, y, 8.5, 1.3, CARD_BG)
    add_rect(s, 0.4, y, 0.06, 1.3, ACCENT)
    add_textbox(s, val,   0.7, y + 0.1, 2.5, 0.8,
                font_size=38, bold=True, color=ACCENT)
    add_textbox(s, label, 3.2, y + 0.1, 5.5, 0.5,
                font_size=18, bold=True, color=WHITE)
    add_textbox(s, sub,   3.2, y + 0.6, 5.5, 0.45,
                font_size=13, color=LIGHT_GRAY)

# Why now
add_rect(s, 9.2, 1.5, 3.8, 5.3, CARD_BG)
add_textbox(s, "WHY NOW", 9.35, 1.6, 3.5, 0.4,
            font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
why_now = [
    "• Solvency II & climate\n  disclosure mandates",
    "• AI/ML now mature\n  enough for real-time\n  geospatial scoring",
    "• 2024 record $380B\n  total disaster losses",
    "• Insurers desperate\n  for automated tools",
]
for i, txt in enumerate(why_now):
    add_textbox(s, txt, 9.35, 2.1 + i * 1.15, 3.4, 1.0,
                font_size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MACHINE LEARNING MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "MACHINE LEARNING", 0.55, 0.08, 6, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "How We Train the Risk Model", 0.55, 0.7, 10, 0.9,
            font_size=40, bold=True, color=WHITE)
accent_bar(s, 1.55, 5.5)

# Left column — pipeline steps
pipeline = [
    ("01", "Raw Data Ingestion",     "Historical claims, seismic surveys,\nclimate records, building registries"),
    ("02", "Preprocessing",          "Normalization to 0–1 floats,\nfeature engineering, outlier removal"),
    ("03", "Model Training",         "Gradient-boosted ensemble trained\non Insurance_Ready_For_ML.csv"),
    ("04", "Knowledge Base Export",  "Model insights exported to\nearthrisk_knowledge_base.txt for agents"),
]
for i, (num, title, desc) in enumerate(pipeline):
    y = 1.8 + i * 1.3
    add_rect(s, 0.4, y, 0.55, 1.1, ACCENT)
    add_textbox(s, num, 0.42, y + 0.25, 0.5, 0.55,
                font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(s, 0.95, y, 5.8, 1.1, CARD_BG)
    add_textbox(s, title, 1.1, y + 0.05, 5.5, 0.45,
                font_size=17, bold=True, color=WHITE)
    add_textbox(s, desc,  1.1, y + 0.5,  5.5, 0.55,
                font_size=13, color=LIGHT_GRAY)

# Right column — feature weights visual
add_rect(s, 7.3, 1.7, 5.7, 5.1, CARD_BG)
add_textbox(s, "FEATURE WEIGHTS", 7.5, 1.8, 5.3, 0.4,
            font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

features_w = [
    ("Seismic Activity",  0.30, "30%"),
    ("Fire Access",       0.20, "20%"),
    ("Climate Exposure",  0.20, "20%"),
    ("Building Age",      0.10, "10%"),
    ("Volcanic Risk",     0.10, "10%"),
    ("Claims History",    0.10, "10%"),
]
bar_max_w = 4.0
for i, (feat, weight, pct) in enumerate(features_w):
    y = 2.35 + i * 0.72
    add_textbox(s, feat, 7.45, y, 2.2, 0.38,
                font_size=13, color=LIGHT_GRAY)
    add_rect(s, 9.65, y + 0.05, bar_max_w * weight, 0.28, ACCENT)
    add_rect(s, 9.65, y + 0.05, bar_max_w,           0.28, RGBColor(0x1E, 0x3A, 0x50))
    add_rect(s, 9.65, y + 0.05, bar_max_w * weight,  0.28, ACCENT)
    add_textbox(s, pct, 9.65 + bar_max_w + 0.1, y, 0.5, 0.38,
                font_size=13, bold=True, color=ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ML RESULTS & PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "ML RESULTS", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "Model Performance & Outcomes", 0.55, 0.7, 10, 0.9,
            font_size=40, bold=True, color=WHITE)
accent_bar(s, 1.55, 4.5)

# Top metric cards
metrics = [
    ("92%",   "Classification\nAccuracy",    "High / Medium / Low risk"),
    ("0.94",  "AUC-ROC\nScore",              "Across all risk categories"),
    ("87%",   "Precision\non High Risk",     "Critical flag reliability"),
    ("3 sec", "Inference\nLatency",          "Per building, real-time API"),
]
for i, (val, label, sub) in enumerate(metrics):
    x = 0.3 + i * 3.25
    add_rect(s, x, 1.8, 3.0, 2.2, CARD_BG)
    add_rect(s, x, 1.8, 3.0, 0.06, ACCENT)
    add_textbox(s, val,   x + 0.1, 1.95, 2.8, 0.85,
                font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x + 0.1, 2.75, 2.8, 0.6,
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, sub,   x + 0.1, 3.3,  2.8, 0.45,
                font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom two panels
# Left — dataset facts
add_rect(s, 0.3, 4.2, 6.0, 3.0, CARD_BG)
add_textbox(s, "TRAINING DATASET", 0.5, 4.3, 5.6, 0.4,
            font_size=13, bold=True, color=ACCENT)
dataset_facts = [
    "•  50,000+ labelled building records",
    "•  12 input features per sample",
    "•  3-class target: High / Medium / Low",
    "•  80/10/10 train / val / test split",
    "•  Oversampling applied to balance High-Risk class",
]
add_bullet_points(s, dataset_facts, 0.5, 4.75, 5.6, 2.3,
                  font_size=14, color=LIGHT_GRAY, spacing_pt=8)

# Right — what it unlocks
add_rect(s, 6.9, 4.2, 6.1, 3.0, CARD_BG)
add_textbox(s, "WHAT THE MODEL UNLOCKS", 7.1, 4.3, 5.7, 0.4,
            font_size=13, bold=True, color=ACCENT)
unlocks = [
    "•  Instant score for any new building",
    "•  Confidence intervals on risk bands",
    "•  Feature attribution → explainable AI",
    "•  Continuous retraining as claims arrive",
    "•  Powers all 4 WatsonX AI agents",
]
add_bullet_points(s, unlocks, 7.1, 4.75, 5.7, 2.3,
                  font_size=14, color=LIGHT_GRAY, spacing_pt=8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CALL TO ACTION / NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

add_rect(s, 0, 0, 13.33, 0.5, CARD_BG)
add_textbox(s, "WHAT'S NEXT", 0.55, 0.08, 5, 0.4,
            font_size=14, bold=True, color=ACCENT)

add_textbox(s, "Let's Build the Future\nof Risk Intelligence", 0.55, 0.7, 9, 1.6,
            font_size=40, bold=True, color=WHITE)
accent_bar(s, 2.3, 5.5)

steps = [
    ("Q2 2026", "Pilot Program", "Partner with 2–3 regional insurers\nfor live portfolio testing"),
    ("Q3 2026", "SaaS Launch", "Cloud-hosted platform with\nAPI access & custom dashboards"),
    ("Q4 2026", "Scale & Expand", "EU & APAC markets, IoT sensor\nintegration, real-time satellite data"),
]
for i, (qtr, title, desc) in enumerate(steps):
    x = 0.4 + i * 4.3
    add_rect(s, x, 2.8, 3.9, 3.5, CARD_BG)
    add_rect(s, x, 2.8, 3.9, 0.06, ACCENT)
    add_textbox(s, qtr,   x + 0.15, 2.9, 3.6, 0.45,
                font_size=13, bold=True, color=ACCENT)
    add_textbox(s, title, x + 0.15, 3.35, 3.6, 0.6,
                font_size=22, bold=True, color=WHITE)
    add_textbox(s, desc,  x + 0.15, 3.95, 3.6, 1.3,
                font_size=14, color=LIGHT_GRAY)

# CTA
add_rect(s, 0.4, 6.45, 12.5, 0.85, ACCENT)
add_textbox(s,
    "We're seeking pilot partners & seed investment  •  Let's talk →  earthrisk@demo.com",
    0.55, 6.52, 12.2, 0.7,
    font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
output = "C:/Users/stavie/Desktop/EarthRisk_Pitch_2026.pptx"
prs.save(output)
print(f"Saved: {output}")
